from __future__ import annotations

import math
import statistics
import unicodedata
from collections import Counter
from typing import Iterable, List, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def round_half(value: float) -> float:
    return round(value * 2) / 2.0


def clamp(value: float, lo: float, hi: float) -> float:
    return max(lo, min(hi, value))


def iter_leaf_shapes(shapes) -> Iterable:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_leaf_shapes(shape.shapes)
        else:
            yield shape


def is_text_shape(shape) -> bool:
    return bool(getattr(shape, "has_text_frame", False)) and bool((shape.text or "").strip())


def shape_rect_pt(shape) -> Tuple[float, float, float, float]:
    return shape.left.pt, shape.top.pt, shape.width.pt, shape.height.pt


def rect_overlap(a, b) -> Tuple[float, float, float]:
    l1, t1, w1, h1 = a
    l2, t2, w2, h2 = b
    r1, b1 = l1 + w1, t1 + h1
    r2, b2 = l2 + w2, t2 + h2
    iw = min(r1, r2) - max(l1, l2)
    ih = min(b1, b2) - max(t1, t2)
    if iw <= 0 or ih <= 0:
        return 0.0, 0.0, 0.0
    return iw, ih, iw * ih


def shape_text_len(shape) -> int:
    return len((shape.text or "").strip())


def resolve_paragraph_font_pt(paragraph, fallback_font_pt: float) -> float:
    if paragraph.font is not None and paragraph.font.size is not None:
        return float(paragraph.font.size.pt)
    for run in paragraph.runs:
        if run.font is not None and run.font.size is not None:
            return float(run.font.size.pt)
    return fallback_font_pt


def resolve_run_font_pt(run, paragraph, fallback_font_pt: float) -> float:
    if run.font is not None and run.font.size is not None:
        return float(run.font.size.pt)
    if paragraph.font is not None and paragraph.font.size is not None:
        return float(paragraph.font.size.pt)
    return fallback_font_pt


def paragraph_line_spacing_pt(paragraph, font_pt: float) -> float:
    line_spacing = paragraph.line_spacing
    if line_spacing is None:
        return font_pt * 1.22
    if hasattr(line_spacing, "pt"):
        return float(line_spacing.pt)
    try:
        value = float(line_spacing)
        if 0.2 <= value <= 6.0:
            return font_pt * value
    except Exception:
        pass
    return font_pt * 1.22


def char_width_factor(ch: str) -> float:
    if ch.isspace():
        return 0.33
    if ch.isdigit():
        return 0.56
    if "A" <= ch <= "Z":
        return 0.62
    if "a" <= ch <= "z":
        return 0.52
    if ch in ".,:;!?锛屻€傦細锛涳紒锛熴€?":
        return 0.45
    if ch in "()[]{}<>锛堬級銆愩€戙€娿€?":
        return 0.58
    if unicodedata.east_asian_width(ch) in ("W", "F"):
        return 1.0
    return 0.62


def estimate_text_width_pt(text: str, font_pt: float) -> float:
    if not text:
        return 0.0
    return sum(char_width_factor(ch) * font_pt for ch in text)


def split_soft_lines(text: str) -> List[str]:
    if not text:
        return [""]
    return text.replace("\r", "\n").replace("\v", "\n").split("\n")


def estimate_paragraph_need_pt(
    paragraph, avail_width_pt: float, fallback_font_pt: float, word_wrap: bool
) -> Tuple[float, int, float]:
    font_pt = resolve_paragraph_font_pt(paragraph, fallback_font_pt)
    line_height = paragraph_line_spacing_pt(paragraph, font_pt)
    text_lines = split_soft_lines(paragraph.text)
    rendered_lines = 0
    for line in text_lines:
        one_line_w = estimate_text_width_pt(line, font_pt)
        if word_wrap and avail_width_pt > 1:
            wrapped = max(1, math.ceil(one_line_w / max(1.0, avail_width_pt)))
        else:
            wrapped = 1
        rendered_lines += wrapped

    space_before = paragraph.space_before.pt if paragraph.space_before else 0.0
    space_after = paragraph.space_after.pt if paragraph.space_after else 0.0
    need = rendered_lines * line_height + space_before + space_after
    return need, rendered_lines, font_pt


def estimate_shape_need_pt(shape, fallback_font_pt: float) -> Tuple[float, int]:
    text_frame = shape.text_frame
    avail_width = shape.width.pt - text_frame.margin_left.pt - text_frame.margin_right.pt
    avail_width = max(8.0, avail_width)
    total = text_frame.margin_top.pt + text_frame.margin_bottom.pt
    total_lines = 0
    for paragraph in text_frame.paragraphs:
        need, lines, _ = estimate_paragraph_need_pt(
            paragraph, avail_width, fallback_font_pt, bool(text_frame.word_wrap)
        )
        total += need
        total_lines += lines
    return total, total_lines


def shape_max_line_width_pt(shape, fallback_font_pt: float) -> float:
    max_width = 0.0
    for paragraph in shape.text_frame.paragraphs:
        text = (paragraph.text or "").strip()
        if not text:
            continue
        font_pt = resolve_paragraph_font_pt(paragraph, fallback_font_pt)
        for soft_line in split_soft_lines(text):
            max_width = max(max_width, estimate_text_width_pt(soft_line, font_pt))
    return max_width


def shape_visual_rect_pt(shape, fallback_font_pt: float) -> Tuple[float, float, float, float]:
    left, top, width, height = shape_rect_pt(shape)
    req_h, _ = estimate_shape_need_pt(shape, fallback_font_pt)
    text_frame = shape.text_frame
    avail_width = max(8.0, width - text_frame.margin_left.pt - text_frame.margin_right.pt)
    max_line_width = shape_max_line_width_pt(shape, fallback_font_pt)

    visual_h = max(height, req_h)
    visual_w = width
    if max_line_width > avail_width:
        visual_w = max(width, max_line_width + text_frame.margin_left.pt + text_frame.margin_right.pt)
    return left, top, visual_w, visual_h


def count_visual_overlaps(
    text_shapes: Sequence, fallback_font_pt: float, min_ratio: float = 0.10
) -> int:
    rects = [shape_visual_rect_pt(shape, fallback_font_pt) for shape in text_shapes]
    count = 0
    for i in range(len(rects)):
        for j in range(i + 1, len(rects)):
            a = rects[i]
            b = rects[j]
            _, _, area = rect_overlap(a, b)
            if area <= 0:
                continue
            ratio = area / max(1.0, min(a[2] * a[3], b[2] * b[3]))
            if ratio >= min_ratio:
                count += 1
    return count


def collect_font_size_counts(prs: Presentation) -> Counter:
    counts = Counter()
    for slide in prs.slides:
        for shape in iter_leaf_shapes(slide.shapes):
            if not is_text_shape(shape):
                continue
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.font and paragraph.font.size:
                    counts[round_half(float(paragraph.font.size.pt))] += 1
                for run in paragraph.runs:
                    if run.font and run.font.size:
                        counts[round_half(float(run.font.size.pt))] += 1
    return counts


def choose_common_font_sizes(counts: Counter) -> List[float]:
    if not counts:
        return [10.0, 11.0, 12.0, 14.0, 16.0]
    total = sum(counts.values())
    min_count = max(2, int(total * 0.03))
    common = [size for size, count in counts.items() if count >= min_count]
    if len(common) < 4:
        common = [size for size, _ in counts.most_common(10)]
    return sorted(set(common))


def choose_fallback_body_font(counts: Counter) -> float:
    samples: List[float] = []
    for size, count in counts.items():
        if size <= 20.0:
            samples.extend([size] * count)
    if not samples:
        for size, count in counts.items():
            samples.extend([size] * count)
    if not samples:
        return 12.0
    return float(statistics.median(samples))


def collect_line_spacing_ratios(prs: Presentation, fallback_font_pt: float) -> List[float]:
    ratios: List[float] = []
    for slide in prs.slides:
        for shape in iter_leaf_shapes(slide.shapes):
            if not is_text_shape(shape):
                continue
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.line_spacing is None:
                    continue
                font_pt = resolve_paragraph_font_pt(paragraph, fallback_font_pt)
                line_spacing_pt = paragraph_line_spacing_pt(paragraph, font_pt)
                ratio = line_spacing_pt / max(0.1, font_pt)
                if 0.9 <= ratio <= 3.0:
                    ratios.append(ratio)
    return ratios


def choose_target_line_ratio(ratios: Sequence[float]) -> float:
    if not ratios:
        return 1.35
    return clamp(float(statistics.median(ratios)), 1.2, 1.5)


def nearest_common_font(size: float, common: Sequence[float], tol: float = 1.2) -> Optional[float]:
    if not common:
        return None
    target = min(common, key=lambda item: abs(item - size))
    if abs(target - size) <= tol:
        return float(target)
    return None

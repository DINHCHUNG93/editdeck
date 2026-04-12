from __future__ import annotations

import re
from typing import Sequence

from pptx.util import Pt

from .metrics import (
    char_width_factor,
    clamp,
    estimate_shape_need_pt,
    estimate_text_width_pt,
    nearest_common_font,
    paragraph_line_spacing_pt,
    resolve_paragraph_font_pt,
    resolve_run_font_pt,
)
from .models import FixTuning


def text_ascii_letter_ratio(text: str) -> float:
    letters = [ch for ch in text if ch.isalpha()]
    if not letters:
        return 0.0
    ascii_letters = sum(1 for ch in letters if "a" <= ch.lower() <= "z")
    return ascii_letters / max(1, len(letters))


def paragraph_has_bullet_hint(paragraph) -> bool:
    xml = paragraph._p.xml
    if "<a:buChar" in xml or "<a:buAutoNum" in xml or "<a:buBlip" in xml:
        return True
    return paragraph.level > 0


def collapse_spaces(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def fix_english_text_artifacts(text: str) -> str:
    text = collapse_spaces(text)
    text = re.sub(
        r"\b([A-Za-z][A-Za-z\-]{2,})\s+and\s+of\s+\1\s+and\b",
        r"\1 and",
        text,
        flags=re.IGNORECASE,
    )
    text = re.sub(
        r"\b(and|or)\s+of\s+([A-Za-z][A-Za-z\-]{2,})\s+(and|or)\b",
        lambda match: match.group(1) if match.group(1).lower() == match.group(3).lower() else match.group(1),
        text,
        flags=re.IGNORECASE,
    )
    return collapse_spaces(text)


def should_flatten_english_paragraphs(shape) -> bool:
    paragraphs = [paragraph for paragraph in shape.text_frame.paragraphs if (paragraph.text or "").strip()]
    if len(paragraphs) < 2 or len(paragraphs) > 8:
        return False
    if any(paragraph_has_bullet_hint(paragraph) for paragraph in paragraphs):
        return False

    texts = [paragraph.text.strip() for paragraph in paragraphs]
    joined = " ".join(texts)
    if len(joined) < 24 or len(joined) > 420:
        return False
    if text_ascii_letter_ratio(joined) < 0.90:
        return False
    if re.search(r"^\s*(\d+[\.\)]|[-•●])\s+", texts[0]):
        return False

    connector_tail = re.compile(
        r"\b(and|or|of|to|in|on|for|from|with|by|as|at|a|an|the|across|via)$",
        re.IGNORECASE,
    )
    has_connector_break = any(connector_tail.search(text) for text in texts[:-1])
    has_weird_spaces = any("   " in text for text in texts)
    has_dup_artifact = bool(
        re.search(
            r"\b(and|or)\s+of\s+([A-Za-z][A-Za-z\-]{2,})\s+(and|or)\b",
            joined,
            flags=re.IGNORECASE,
        )
    )
    return has_weird_spaces or has_dup_artifact or (has_connector_break and len(paragraphs) >= 4)


def flatten_english_hard_breaks(shape) -> int:
    if not should_flatten_english_paragraphs(shape):
        return 0

    paragraphs = [paragraph for paragraph in shape.text_frame.paragraphs if (paragraph.text or "").strip()]
    merged = fix_english_text_artifacts(" ".join(paragraph.text.strip() for paragraph in paragraphs))
    if not merged:
        return 0

    first = shape.text_frame.paragraphs[0]
    old = collapse_spaces(" ".join(paragraph.text.strip() for paragraph in paragraphs))
    if merged == old and len(paragraphs) == 1:
        return 0

    first.text = merged
    first.level = 0
    for paragraph in list(shape.text_frame.paragraphs[1:]):
        paragraph._p.getparent().remove(paragraph._p)
    return 1


def is_chinese_char(ch: str) -> bool:
    return "\u4e00" <= ch <= "\u9fff"


def fix_short_chinese_breaks(shape) -> int:
    paragraphs = [paragraph for paragraph in shape.text_frame.paragraphs if (paragraph.text or "").strip()]
    if len(paragraphs) < 2 or len(paragraphs) > 4:
        return 0
    if any(paragraph_has_bullet_hint(paragraph) for paragraph in paragraphs):
        return 0

    full_text = "\n".join((paragraph.text or "").strip() for paragraph in paragraphs)
    pure = full_text.replace("\n", "")
    if len(pure) < 6 or len(pure) > 28:
        return 0
    if text_ascii_letter_ratio(pure) > 0.10:
        return 0
    if shape.height.pt > 48.0:
        return 0

    connectors = set("与和及并或在于对将把给向从到为由并且而且及其")
    parts = [(paragraph.text or "").strip() for paragraph in paragraphs]
    merged = parts[0]
    changed = False
    for nxt in parts[1:]:
        prev_char = merged[-1] if merged else ""
        next_char = nxt[0] if nxt else ""
        should_merge = (
            bool(prev_char)
            and bool(next_char)
            and is_chinese_char(prev_char)
            and is_chinese_char(next_char)
            and prev_char not in connectors
            and len(nxt) <= 4
        )
        if should_merge:
            merged += nxt
            changed = True
        else:
            merged += "\n" + nxt

    if not changed:
        return 0

    first = shape.text_frame.paragraphs[0]
    first.text = merged
    first.level = 0
    for paragraph in list(shape.text_frame.paragraphs[1:]):
        paragraph._p.getparent().remove(paragraph._p)
    return 1


def snap_font_outliers(shape, common_sizes: Sequence[float]) -> int:
    changed = 0
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.font and paragraph.font.size:
            current = float(paragraph.font.size.pt)
            target = nearest_common_font(current, common_sizes)
            if target is not None and abs(target - current) >= 0.45:
                paragraph.font.size = Pt(target)
                changed += 1
        for run in paragraph.runs:
            if run.font and run.font.size:
                current = float(run.font.size.pt)
                target = nearest_common_font(current, common_sizes)
                if target is not None and abs(target - current) >= 0.45:
                    run.font.size = Pt(target)
                    changed += 1
    return changed


def normalize_line_spacing(shape, fallback_font_pt: float, target_ratio: float) -> int:
    changed = 0
    for paragraph in shape.text_frame.paragraphs:
        if not (paragraph.text or "").strip():
            continue
        font_pt = resolve_paragraph_font_pt(paragraph, fallback_font_pt)
        if paragraph.line_spacing is None:
            continue
        line_spacing_pt = paragraph_line_spacing_pt(paragraph, font_pt)
        ratio = line_spacing_pt / max(0.1, font_pt)
        if ratio > 1.75:
            paragraph.line_spacing = Pt(round(font_pt * target_ratio, 1))
            changed += 1
    return changed


def tighten_text_margins(shape, force: bool = False) -> int:
    text_frame = shape.text_frame
    text = (shape.text or "").strip()
    if not text:
        return 0

    height = shape.height.pt
    width = shape.width.pt
    paragraphs = [paragraph for paragraph in text_frame.paragraphs if (paragraph.text or "").strip()]
    compact = height <= 22.0 or (height <= 30.0 and len(paragraphs) <= 1) or (width <= 90.0 and height <= 42.0)
    if not compact and not force:
        return 0

    target_lr = 0.8 if width <= 140.0 else 1.2
    target_tb = 0.4 if height <= 22.0 else 0.8
    changed = 0

    if text_frame.margin_left.pt > target_lr + 0.05:
        text_frame.margin_left = Pt(target_lr)
        changed += 1
    if text_frame.margin_right.pt > target_lr + 0.05:
        text_frame.margin_right = Pt(target_lr)
        changed += 1
    if text_frame.margin_top.pt > target_tb + 0.05:
        text_frame.margin_top = Pt(target_tb)
        changed += 1
    if text_frame.margin_bottom.pt > target_tb + 0.05:
        text_frame.margin_bottom = Pt(target_tb)
        changed += 1
    return changed


def scale_text_shape_fonts(shape, scale: float, fallback_font_pt: float, min_font_pt: float = 8.0) -> int:
    scale = clamp(scale, 0.5, 1.5)
    changed = 0
    for paragraph in shape.text_frame.paragraphs:
        base_p = resolve_paragraph_font_pt(paragraph, fallback_font_pt)
        target_p = max(min_font_pt, round(base_p * scale, 1))
        if paragraph.font.size is not None and abs(float(paragraph.font.size.pt) - target_p) >= 0.11:
            paragraph.font.size = Pt(target_p)
            changed += 1
        run_has_explicit = any(run.font and run.font.size for run in paragraph.runs)
        if paragraph.font.size is None and not run_has_explicit:
            paragraph.font.size = Pt(target_p)
            changed += 1
        for run in paragraph.runs:
            base_r = resolve_run_font_pt(run, paragraph, base_p)
            target_r = max(min_font_pt, round(base_r * scale, 1))
            if run.font.size is not None and abs(float(run.font.size.pt) - target_r) >= 0.11:
                run.font.size = Pt(target_r)
                changed += 1
    return changed


def weighted_text_len(text: str) -> float:
    return sum(char_width_factor(ch) for ch in text)


def enforce_single_line_like_text(shape, fallback_font_pt: float, tuning: FixTuning) -> int:
    text = (shape.text or "").strip()
    if not text:
        return 0
    if "\n" in text or "\r" in text or "\v" in text:
        return 0
    if len(shape.text_frame.paragraphs) != 1:
        return 0

    width_units = weighted_text_len(text)
    if width_units > tuning.single_width_units_limit or len(text) > tuning.single_text_len_limit:
        return 0

    text_frame = shape.text_frame
    paragraph = text_frame.paragraphs[0]
    font_pt = resolve_paragraph_font_pt(paragraph, fallback_font_pt)
    avail_width = shape.width.pt - text_frame.margin_left.pt - text_frame.margin_right.pt
    avail_width = max(8.0, avail_width)
    current_width = estimate_text_width_pt(text, font_pt)
    if current_width <= avail_width * 1.02:
        return 0
    if current_width >= avail_width * tuning.single_overflow_ratio_max:
        return 0

    target_scale = clamp((avail_width * 0.98) / max(1.0, current_width), tuning.single_scale_floor, 0.99)
    changed = scale_text_shape_fonts(shape, target_scale, fallback_font_pt, min_font_pt=tuning.min_font_pt)
    if paragraph.line_spacing is not None:
        new_font = resolve_paragraph_font_pt(paragraph, fallback_font_pt)
        paragraph.line_spacing = Pt(round(new_font * 1.2, 1))
        changed += 1
    return changed


def fit_shape_overflow(
    shape,
    slide_height_pt: float,
    fallback_font_pt: float,
    target_line_ratio: float,
    tuning: FixTuning,
    allow_expand_box: bool = False,
) -> int:
    changed = 0
    req, _ = estimate_shape_need_pt(shape, fallback_font_pt)
    box = shape.height.pt
    if req <= box * tuning.overflow_ignore_ratio:
        return 0

    changed += tighten_text_margins(shape, force=req > box * 1.18)
    req, _ = estimate_shape_need_pt(shape, fallback_font_pt)
    box = shape.height.pt
    if req <= box * tuning.overflow_ignore_ratio:
        return changed

    changed += normalize_line_spacing(shape, fallback_font_pt, target_line_ratio)
    req, _ = estimate_shape_need_pt(shape, fallback_font_pt)
    box = shape.height.pt

    if req > box * tuning.overflow_scale_trigger_ratio:
        scale = clamp((box * 1.03) / max(1.0, req), tuning.overflow_scale_floor, 0.98)
        if scale < 0.995:
            changed += scale_text_shape_fonts(shape, scale, fallback_font_pt, min_font_pt=tuning.min_font_pt)

    if req > box * 1.25 and allow_expand_box:
        extra_room = slide_height_pt - (shape.top.pt + box)
        need = req - box + 2.0
        if extra_room > 2.0:
            add = min(extra_room, need)
            if add > 1.0:
                shape.height = Pt(box + add)
                changed += 1

    return changed

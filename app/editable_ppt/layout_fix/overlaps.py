from __future__ import annotations

from typing import Sequence

from pptx.util import Pt

from .metrics import (
    clamp,
    estimate_shape_need_pt,
    is_text_shape,
    iter_leaf_shapes,
    rect_overlap,
    shape_rect_pt,
    shape_text_len,
    shape_visual_rect_pt,
)
from .models import FixTuning
from .text_ops import normalize_line_spacing, scale_text_shape_fonts


def resolve_visual_text_overlaps(
    slide,
    fallback_font_pt: float,
    target_line_ratio: float,
    tuning: FixTuning,
    rounds: int = 3,
) -> int:
    changed = 0
    text_shapes = [shape for shape in iter_leaf_shapes(slide.shapes) if is_text_shape(shape)]
    if len(text_shapes) < 2:
        return 0

    for _ in range(rounds):
        req_map = {id(shape): estimate_shape_need_pt(shape, fallback_font_pt)[0] for shape in text_shapes}
        constraints = {}

        for i in range(len(text_shapes)):
            for j in range(i + 1, len(text_shapes)):
                a = text_shapes[i]
                b = text_shapes[j]
                rect_a = shape_visual_rect_pt(a, fallback_font_pt)
                rect_b = shape_visual_rect_pt(b, fallback_font_pt)
                iw, ih, area = rect_overlap(rect_a, rect_b)
                if area <= 0 or iw < 4.0 or ih < 3.0:
                    continue

                if rect_a[1] <= rect_b[1]:
                    upper, lower = a, b
                else:
                    upper, lower = b, a

                max_upper_h = lower.top.pt - upper.top.pt - 2.0
                if max_upper_h <= 8.0:
                    continue
                upper_req = req_map[id(upper)]
                if upper_req > max_upper_h * 1.01:
                    prev = constraints.get(id(upper), float("inf"))
                    constraints[id(upper)] = min(prev, max_upper_h)

        if not constraints:
            break

        changed_this_round = 0
        for shape in text_shapes:
            if id(shape) not in constraints:
                continue
            target_req = constraints[id(shape)]
            req = estimate_shape_need_pt(shape, fallback_font_pt)[0]
            if req <= target_req * 1.01:
                continue
            normalize_line_spacing(shape, fallback_font_pt, target_line_ratio)
            req = estimate_shape_need_pt(shape, fallback_font_pt)[0]
            if req <= target_req * 1.01:
                continue

            scale = clamp((target_req * 0.99) / max(1.0, req), tuning.overflow_scale_floor, 0.98)
            if scale >= 0.995:
                continue
            changed_this_round += scale_text_shape_fonts(
                shape, scale, fallback_font_pt, min_font_pt=tuning.min_font_pt
            )

        changed += changed_this_round
        if changed_this_round == 0:
            break

    return changed


def resolve_compact_box_overlaps(
    slide, slide_w_pt: float, slide_h_pt: float, min_ratio: float = 0.55
) -> int:
    changed = 0
    text_shapes = [shape for shape in iter_leaf_shapes(slide.shapes) if is_text_shape(shape)]
    if len(text_shapes) < 2:
        return 0

    for _ in range(4):
        best = None
        for i in range(len(text_shapes)):
            for j in range(i + 1, len(text_shapes)):
                a = text_shapes[i]
                b = text_shapes[j]
                rect_a = shape_rect_pt(a)
                rect_b = shape_rect_pt(b)
                iw, ih, area = rect_overlap(rect_a, rect_b)
                if area <= 0:
                    continue
                ratio = area / max(1.0, min(rect_a[2] * rect_a[3], rect_b[2] * rect_b[3]))
                if ratio < min_ratio:
                    continue
                if best is None or ratio > best[0]:
                    best = (ratio, iw, ih, a, b)

        if best is None:
            break

        _, iw, ih, a, b = best
        a_len = shape_text_len(a)
        b_len = shape_text_len(b)
        a_area = a.width.pt * a.height.pt
        b_area = b.width.pt * b.height.pt
        if a_area <= b_area and a_len <= b_len:
            mover = a
        elif b_area <= a_area and b_len <= a_len:
            mover = b
        else:
            mover = a if a_area <= b_area else b
        anchor = b if mover is a else a

        if mover.height.pt > 30.0 and shape_text_len(mover) > 20:
            break

        mover_left, mover_top, mover_width, mover_height = shape_rect_pt(mover)
        anchor_left, anchor_top, anchor_width, anchor_height = shape_rect_pt(anchor)
        old_score = overlap_score_with_others(mover, text_shapes)

        x_candidates = {
            mover_left,
            anchor_left + anchor_width + 4.0,
            anchor_left - mover_width - 4.0,
            mover_left + iw + 4.0,
            mover_left - iw - 4.0,
            slide_w_pt - mover_width - 4.0,
            4.0,
        }
        y_candidates = {
            mover_top,
            anchor_top + anchor_height + 2.0,
            anchor_top - mover_height - 2.0,
            mover_top + ih + 2.0,
            mover_top - ih - 2.0,
            mover_top + 18.0,
            mover_top - 18.0,
            4.0,
        }
        for dx in range(-160, 161, 16):
            x_candidates.add(mover_left + dx)
        for dy in range(-120, 121, 6):
            y_candidates.add(mover_top + dy)
        candidates = [(x, y) for x in x_candidates for y in y_candidates]

        best_pos = None
        best_score = old_score
        for new_left, new_top in candidates:
            new_left = clamp(new_left, 0.0, max(0.0, slide_w_pt - mover_width))
            new_top = clamp(new_top, 0.0, max(0.0, slide_h_pt - mover_height))
            if abs(new_left - mover_left) < 0.05 and abs(new_top - mover_top) < 0.05:
                continue
            mover.left = Pt(new_left)
            mover.top = Pt(new_top)
            overlap_score = overlap_score_with_others(mover, text_shapes)
            move_penalty = abs(new_left - mover_left) + abs(new_top - mover_top)
            score = overlap_score + move_penalty * 1.2
            if score + 5.0 < best_score:
                best_score = score
                best_pos = (new_left, new_top)

        mover.left = Pt(mover_left)
        mover.top = Pt(mover_top)
        if best_pos is None:
            break
        mover.left = Pt(best_pos[0])
        mover.top = Pt(best_pos[1])
        changed += 1

    return changed


def compress_bottom_region_to_fit(
    slide,
    slide_h_pt: float,
    threshold_ratio: float = 0.55,
    overflow_trigger_pt: float = 10.0,
    margin_pt: float = 2.0,
) -> int:
    leaf_shapes = list(iter_leaf_shapes(slide.shapes))
    if not leaf_shapes:
        return 0

    max_bottom = 0.0
    for shape in leaf_shapes:
        try:
            bottom = shape.top.pt + shape.height.pt
        except Exception:
            continue
        if bottom > max_bottom:
            max_bottom = bottom

    overflow = max_bottom - slide_h_pt
    if overflow <= overflow_trigger_pt:
        return 0

    threshold = slide_h_pt * threshold_ratio
    span = max_bottom - threshold
    target_span = slide_h_pt - threshold - margin_pt
    if span <= 1.0 or target_span <= 1.0:
        return 0
    scale = min(1.0, target_span / span)
    if scale >= 0.995:
        return 0

    changed = 0
    for shape in leaf_shapes:
        try:
            top = shape.top.pt
            height = shape.height.pt
        except Exception:
            continue
        if top >= threshold or top + height > slide_h_pt:
            new_top = threshold + (top - threshold) * scale
            new_height = max(0.0, height * scale) if height <= 2.0 else max(6.0, height * scale)
            shape.top = Pt(max(0.0, new_top))
            shape.height = Pt(new_height)
            changed += 1

    return changed


def clamp_shape_in_slide(shape, slide_w_pt: float, slide_h_pt: float) -> int:
    changed = 0
    left = shape.left.pt
    top = shape.top.pt
    width = shape.width.pt
    height = shape.height.pt

    if width > slide_w_pt:
        shape.width = Pt(slide_w_pt)
        left = 0.0
        changed += 1
    if height > slide_h_pt:
        shape.height = Pt(slide_h_pt)
        top = 0.0
        changed += 1

    if left < 0:
        shape.left = Pt(0)
        changed += 1
    if top < 0:
        shape.top = Pt(0)
        changed += 1

    right = shape.left.pt + shape.width.pt
    bottom = shape.top.pt + shape.height.pt
    if right > slide_w_pt:
        shape.left = Pt(max(0.0, slide_w_pt - shape.width.pt))
        changed += 1
    if bottom > slide_h_pt:
        shape.top = Pt(max(0.0, slide_h_pt - shape.height.pt))
        changed += 1

    return changed


def choose_shape_to_move(shape_a, shape_b, slide_h_pt: float):
    a_placeholder = bool(getattr(shape_a, "is_placeholder", False))
    b_placeholder = bool(getattr(shape_b, "is_placeholder", False))
    if a_placeholder and not b_placeholder:
        return shape_b
    if b_placeholder and not a_placeholder:
        return shape_a

    a_rect = shape_rect_pt(shape_a)
    b_rect = shape_rect_pt(shape_b)
    a_area = a_rect[2] * a_rect[3]
    b_area = b_rect[2] * b_rect[3]
    a_len = shape_text_len(shape_a)
    b_len = shape_text_len(shape_b)

    a_is_title = a_rect[1] < slide_h_pt * 0.2 and a_rect[3] < slide_h_pt * 0.2
    b_is_title = b_rect[1] < slide_h_pt * 0.2 and b_rect[3] < slide_h_pt * 0.2
    if a_is_title and not b_is_title:
        return shape_b
    if b_is_title and not a_is_title:
        return shape_a

    if a_area < b_area * 0.9:
        return shape_a
    if b_area < a_area * 0.9:
        return shape_b
    return shape_a if a_len <= b_len else shape_b


def overlap_score_with_others(shape, others) -> float:
    base = shape_rect_pt(shape)
    total = 0.0
    for other in others:
        if other is shape:
            continue
        _, _, area = rect_overlap(base, shape_rect_pt(other))
        total += area
    return total


def count_major_overlaps(text_shapes: Sequence, min_ratio: float = 0.22) -> int:
    count = 0
    for i in range(len(text_shapes)):
        for j in range(i + 1, len(text_shapes)):
            a = shape_rect_pt(text_shapes[i])
            b = shape_rect_pt(text_shapes[j])
            _, _, area = rect_overlap(a, b)
            if area <= 0:
                continue
            ratio = area / max(1.0, min(a[2] * a[3], b[2] * b[3]))
            if ratio >= min_ratio:
                count += 1
    return count


def resolve_major_overlaps(
    slide, slide_w_pt: float, slide_h_pt: float, min_ratio_move: float = 0.35
) -> int:
    changed = 0
    text_shapes = [shape for shape in iter_leaf_shapes(slide.shapes) if is_text_shape(shape)]
    if len(text_shapes) < 2:
        return 0

    for _ in range(4):
        moved_this_round = 0
        pairs = []
        for i in range(len(text_shapes)):
            for j in range(i + 1, len(text_shapes)):
                a = text_shapes[i]
                b = text_shapes[j]
                rect_a = shape_rect_pt(a)
                rect_b = shape_rect_pt(b)
                iw, ih, area = rect_overlap(rect_a, rect_b)
                if area <= 0:
                    continue
                ratio = area / max(1.0, min(rect_a[2] * rect_a[3], rect_b[2] * rect_b[3]))
                if ratio >= min_ratio_move:
                    pairs.append((ratio, iw, ih, a, b))

        if not pairs:
            break
        pairs.sort(key=lambda item: item[0], reverse=True)

        for _, iw, ih, a, b in pairs:
            mover = choose_shape_to_move(a, b, slide_h_pt)
            left, top, width, height = shape_rect_pt(mover)
            old_score = overlap_score_with_others(mover, text_shapes)
            candidates = [
                (left, top + ih + 6.0),
                (left + iw + 6.0, top),
                (left, top - ih - 6.0),
                (left - iw - 6.0, top),
            ]
            best_pos = None
            best_score = old_score

            for new_left, new_top in candidates:
                new_left = clamp(new_left, 0.0, max(0.0, slide_w_pt - width))
                new_top = clamp(new_top, 0.0, max(0.0, slide_h_pt - height))
                if abs(new_left - left) < 0.05 and abs(new_top - top) < 0.05:
                    continue
                mover.left = Pt(new_left)
                mover.top = Pt(new_top)
                score = overlap_score_with_others(mover, text_shapes)
                if score + 8.0 < best_score:
                    best_score = score
                    best_pos = (new_left, new_top)

            mover.left = Pt(left)
            mover.top = Pt(top)
            if best_pos is not None:
                mover.left = Pt(best_pos[0])
                mover.top = Pt(best_pos[1])
                changed += 1
                moved_this_round += 1

        if moved_this_round == 0:
            break

    return changed

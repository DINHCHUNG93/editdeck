from __future__ import annotations

import json
from collections import Counter
from pathlib import Path
from typing import Tuple

from pptx import Presentation

from .metrics import (
    choose_common_font_sizes,
    choose_fallback_body_font,
    choose_target_line_ratio,
    collect_font_size_counts,
    collect_line_spacing_ratios,
    count_visual_overlaps,
    estimate_shape_need_pt,
    is_text_shape,
    iter_leaf_shapes,
)
from .models import AnalyzeResult, FixTuning, tuning_from_mode
from .overlaps import (
    clamp_shape_in_slide,
    compress_bottom_region_to_fit,
    count_major_overlaps,
    resolve_compact_box_overlaps,
    resolve_major_overlaps,
    resolve_visual_text_overlaps,
)
from .text_ops import (
    enforce_single_line_like_text,
    fit_shape_overflow,
    fix_short_chinese_breaks,
    flatten_english_hard_breaks,
    normalize_line_spacing,
    snap_font_outliers,
    tighten_text_margins,
)


def analyze_presentation(prs: Presentation, fallback_font_pt: float) -> AnalyzeResult:
    result = AnalyzeResult()
    for slide in prs.slides:
        text_shapes = [shape for shape in iter_leaf_shapes(slide.shapes) if is_text_shape(shape)]
        result.text_shapes += len(text_shapes)
        for shape in text_shapes:
            req, _ = estimate_shape_need_pt(shape, fallback_font_pt)
            if req > shape.height.pt * 1.03:
                result.overflow_like += 1

        result.major_overlap_pairs += count_major_overlaps(text_shapes, min_ratio=0.22)
        result.visual_overlap_pairs += count_visual_overlaps(text_shapes, fallback_font_pt, min_ratio=0.10)
    return result


def process_presentation(
    path: Path,
    out_path: Path,
    tuning: FixTuning,
    fix_overlap: bool = False,
    allow_expand_box: bool = False,
    clamp_boundary: bool = False,
) -> Tuple[AnalyzeResult, AnalyzeResult, Counter]:
    prs = Presentation(path)
    font_counts = collect_font_size_counts(prs)
    common_sizes = choose_common_font_sizes(font_counts)
    fallback_font_pt = choose_fallback_body_font(font_counts)
    line_ratios = collect_line_spacing_ratios(prs, fallback_font_pt)
    target_line_ratio = choose_target_line_ratio(line_ratios)

    before = analyze_presentation(prs, fallback_font_pt)
    stats = Counter()

    for slide in prs.slides:
        slide_w = prs.slide_width.pt
        slide_h = prs.slide_height.pt
        stats["bottom_compress_fix"] += compress_bottom_region_to_fit(
            slide, slide_h, threshold_ratio=0.55, overflow_trigger_pt=10.0, margin_pt=2.0
        )
        text_shapes = [shape for shape in iter_leaf_shapes(slide.shapes) if is_text_shape(shape)]
        for shape in text_shapes:
            stats["english_break_fix"] += flatten_english_hard_breaks(shape)
            stats["cn_break_fix"] += fix_short_chinese_breaks(shape)
            req_now, _ = estimate_shape_need_pt(shape, fallback_font_pt)
            if req_now > shape.height.pt * 1.02 or shape.height.pt <= 20.0:
                stats["margin_fix"] += tighten_text_margins(shape, force=req_now > shape.height.pt * 1.18)
            stats["font_snap"] += snap_font_outliers(shape, common_sizes)
            stats["line_spacing"] += normalize_line_spacing(shape, fallback_font_pt, target_line_ratio)
            stats["single_line_fix"] += enforce_single_line_like_text(shape, fallback_font_pt, tuning)
            stats["overflow_fix"] += fit_shape_overflow(
                shape,
                slide_h,
                fallback_font_pt,
                target_line_ratio,
                tuning,
                allow_expand_box=allow_expand_box,
            )
            if clamp_boundary:
                stats["clamp"] += clamp_shape_in_slide(shape, slide_w, slide_h)

        stats["visual_overlap_fix"] += resolve_visual_text_overlaps(
            slide, fallback_font_pt, target_line_ratio, tuning
        )
        if fix_overlap:
            stats["compact_overlap_fix"] += resolve_compact_box_overlaps(
                slide, slide_w, slide_h, min_ratio=0.55
            )

        if fix_overlap:
            old_positions = {id(shape): (shape.left, shape.top) for shape in text_shapes}
            before_overlap = count_major_overlaps(text_shapes, min_ratio=0.22)
            moved = resolve_major_overlaps(slide, slide_w, slide_h, min_ratio_move=0.35)
            after_overlap = count_major_overlaps(text_shapes, min_ratio=0.22)
            if after_overlap <= before_overlap:
                stats["overlap_fix"] += moved
            else:
                for shape in text_shapes:
                    left, top = old_positions[id(shape)]
                    shape.left = left
                    shape.top = top
                stats["overlap_reverted"] += 1

        if clamp_boundary:
            for shape in text_shapes:
                stats["clamp"] += clamp_shape_in_slide(shape, slide_w, slide_h)

    after = analyze_presentation(prs, fallback_font_pt)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return before, after, stats


def apply_layout_fix(
    final_path: Path,
    output_root: Path,
    mode: str = "balanced",
    *,
    fix_overlap: bool = False,
    allow_expand_box: bool = False,
    clamp_boundary: bool = False,
) -> str:
    tuning = tuning_from_mode(mode)
    temp_path = output_root / f"{final_path.stem}.layout_fix_tmp{final_path.suffix}"
    report_path = output_root / "layout_fix_report.json"
    if temp_path.exists():
        temp_path.unlink()

    try:
        before, after, stats = process_presentation(
            final_path,
            temp_path,
            tuning=tuning,
            fix_overlap=fix_overlap,
            allow_expand_box=allow_expand_box,
            clamp_boundary=clamp_boundary,
        )
        temp_path.replace(final_path)
    except Exception:
        if temp_path.exists():
            temp_path.unlink()
        raise

    report_path.write_text(
        json.dumps(
            {
                "source_pptx_path": str(final_path),
                "mode": mode,
                "options": {
                    "fix_overlap": fix_overlap,
                    "allow_expand_box": allow_expand_box,
                    "clamp_boundary": clamp_boundary,
                },
                "before": {
                    "text_shapes": before.text_shapes,
                    "overflow_like": before.overflow_like,
                    "major_overlap_pairs": before.major_overlap_pairs,
                    "visual_overlap_pairs": before.visual_overlap_pairs,
                },
                "after": {
                    "text_shapes": after.text_shapes,
                    "overflow_like": after.overflow_like,
                    "major_overlap_pairs": after.major_overlap_pairs,
                    "visual_overlap_pairs": after.visual_overlap_pairs,
                },
                "stats": {key: int(value) for key, value in stats.items()},
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    return (
        "Layout fix applied "
        f"(overflow_like {before.overflow_like}->{after.overflow_like}, "
        f"visual_overlap_pairs {before.visual_overlap_pairs}->{after.visual_overlap_pairs})."
    )
